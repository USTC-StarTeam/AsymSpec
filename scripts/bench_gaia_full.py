"""GAIA full validation bench — all 165 samples (L1+L2+L3).

Extends bench_gaia_vl.py to cover:
  * No-file / web-only questions  (127 samples)
  * Spreadsheet: .xlsx / .xls     (13 samples)
  * PDF                            ( 3 samples)
  * Audio: .mp3                    ( 3 samples)  — transcript placeholder
  * Archive: .zip                  ( 2 samples)
  * Other docs: .pptx .docx .txt .csv .py .jsonld .pdb (7 samples)
  * Image: .png / .jpg             (10 samples)  — VL drafter path

Cells (--cfg):
  floor   : Qwen3-32B alone, verifier sees [truncated file content +] question
  aug     : Qwen3-4B alone, drafter sees full file content + question
  ss      : SpecSteer (Qwen3-32B verifier + Qwen3-4B text drafter, asymmetric context)
  ss_vl   : SpecSteer (Qwen3-32B verifier + Qwen3-VL-2B drafter) — image samples
            get VL drafter; non-image samples fall back to text drafter

For SpecSteer modes, verifier receives truncated content (--trunc_tokens, default 2000)
while drafter receives the full extracted text.

Image samples in cfg=floor/aug/ss receive the VL-generated caption (if captions.json
from bench_gaia_vl.py exists); otherwise a placeholder is used.

Results are saved per cfg with a combined summary suitable for comparison with
bench_gaia_vl.py outputs.
"""
from __future__ import annotations
import argparse, csv, io, json, os, re, sys, time, zipfile
from pathlib import Path
from typing import Optional

import pandas as pd
from PIL import Image

os.environ.setdefault("VLLM_ENABLE_V1_MULTIPROCESSING", "0")

_REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), os.pardir))
if _REPO not in sys.path:
    sys.path.insert(0, _REPO)

from paths import get_hf_token, gaia_validation_dir  # noqa: E402

if not os.environ.get("HF_TOKEN"):
    os.environ["HF_TOKEN"] = get_hf_token() or ""

TEXT_PATH   = "Qwen/Qwen3-32B"
DRAFTER_PATH = "Qwen/Qwen3-4B"
VL_PATH     = "Qwen/Qwen3-VL-2B-Instruct"

OUT_DIR = f"{_REPO}/outputs/gaia_full"
CAPTION_PATH = f"{_REPO}/outputs/gaia_vl/captions.json"

COT_SUFFIX = "\nThink step by step, then provide the final answer after 'Answer:'."

PROMPT_WITH_CTX = (
    "Use the following file content to answer the question.\n\n"
    "{ctx}\n\n"
    "Question: {question}" + COT_SUFFIX
)
PROMPT_NO_CTX = "Question: {question}" + COT_SUFFIX


# ─── file content extractors ──────────────────────────────────────────────────

def _extract_xlsx(fp: str) -> str:
    try:
        sheets = pd.read_excel(fp, sheet_name=None, dtype=str)
        parts = []
        for name, df in sheets.items():
            header = f"[Sheet: {name}]" if len(sheets) > 1 else ""
            parts.append((header + "\n" + df.to_string(index=False)).strip())
        return "\n\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction failed: {e}]"


def _extract_xls(fp: str) -> str:
    return _extract_xlsx(fp)


def _extract_pdf(fp: str) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(fp)
        return "\n\n".join(
            p.extract_text() or "" for p in reader.pages
        ).strip()
    except Exception as e:
        return f"[PDF extraction failed: {e}]"


def _extract_docx(fp: str) -> str:
    try:
        from docx import Document
        doc = Document(fp)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX extraction failed: {e}]"


def _extract_pptx(fp: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(fp)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = [sh.text for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[PPTX extraction failed: {e}]"


def _extract_txt(fp: str) -> str:
    try:
        return Path(fp).read_text(errors="replace").strip()
    except Exception as e:
        return f"[TXT read failed: {e}]"


def _extract_csv(fp: str) -> str:
    try:
        df = pd.read_csv(fp, dtype=str)
        return df.to_string(index=False)
    except Exception as e:
        return f"[CSV extraction failed: {e}]"


def _extract_py(fp: str) -> str:
    return "```python\n" + _extract_txt(fp) + "\n```"


def _extract_jsonld(fp: str) -> str:
    try:
        data = json.loads(Path(fp).read_text())
        return json.dumps(data, indent=2)
    except Exception as e:
        return f"[JSONLD read failed: {e}]"


def _extract_pdb(fp: str) -> str:
    text = _extract_txt(fp)
    # PDB files can be large; keep ATOM/HETATM records + HEADER/REMARK lines
    lines = text.splitlines()
    kept = [l for l in lines if l.startswith(("HEADER", "TITLE", "COMPND",
             "REMARK", "SEQRES", "ATOM", "HETATM", "END"))]
    return "\n".join(kept[:500])  # cap at 500 lines


def _extract_mp3(fp: str) -> str:
    # Whisper ASR — only if openai-whisper available; else placeholder
    try:
        import whisper
        model = whisper.load_model("base")
        result = model.transcribe(fp, fp16=False, language="en")
        return f"[AUDIO TRANSCRIPT]\n{result['text'].strip()}"
    except ImportError:
        return "[AUDIO FILE: whisper not available — transcript omitted]"
    except Exception as e:
        return f"[AUDIO transcription failed: {e}]"


def _extract_zip(fp: str) -> str:
    """Extract all recognisable text from the zip archive."""
    parts: list[str] = []
    try:
        with zipfile.ZipFile(fp) as zf:
            for name in zf.namelist():
                ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
                if ext in ("__macosx", "ds_store"):
                    continue
                try:
                    data = zf.read(name)
                    tmpf = f"/tmp/_gaia_zip_{os.getpid()}_{name.replace('/','_')}"
                    with open(tmpf, "wb") as f:
                        f.write(data)
                    extractor = {
                        "xlsx": _extract_xlsx, "xls": _extract_xls,
                        "pdf": _extract_pdf,   "docx": _extract_docx,
                        "pptx": _extract_pptx, "txt": _extract_txt,
                        "csv": _extract_csv,   "py": _extract_py,
                        "json": _extract_jsonld, "xml": _extract_txt,
                    }.get(ext, _extract_txt)
                    text = extractor(tmpf)
                    parts.append(f"[File in archive: {name}]\n{text}")
                    os.unlink(tmpf)
                except Exception as e:
                    parts.append(f"[File in archive: {name} — failed: {e}]")
    except Exception as e:
        return f"[ZIP extraction failed: {e}]"
    return "\n\n".join(parts)


_EXTRACTORS = {
    "xlsx": _extract_xlsx, "xls": _extract_xls,
    "pdf": _extract_pdf,   "docx": _extract_docx,
    "pptx": _extract_pptx, "txt": _extract_txt,
    "csv": _extract_csv,   "py": _extract_py,
    "jsonld": _extract_jsonld, "pdb": _extract_pdb,
    "mp3": _extract_mp3,   "zip": _extract_zip,
    "png": None, "jpg": None, "jpeg": None,  # handled by VL path
}


# ─── sample loading ────────────────────────────────────────────────────────────

def load_all_samples() -> list[dict]:
    base = str(gaia_validation_dir())
    rows = []
    for L in [1, 2, 3]:
        df = pd.read_parquet(f"{base}/metadata.level{L}.parquet")
        df["level"] = L
        rows.append(df)
    df = pd.concat(rows, ignore_index=True)
    df["ext"] = df["file_name"].fillna("").str.extract(
        r"\.([a-zA-Z0-9]+)$")[0].str.lower()

    captions = json.load(open(CAPTION_PATH)) if os.path.exists(CAPTION_PATH) else {}

    samples = []
    skipped = 0
    for _, r in df.iterrows():
        ext = r["ext"] if isinstance(r["ext"], str) else ""
        fn  = r["file_name"] if isinstance(r["file_name"], str) else ""
        fp  = os.path.join(base, fn) if fn else ""

        is_image = ext in ("png", "jpg", "jpeg")
        file_text = ""
        image_obj: Optional[Image.Image] = None

        if fn and fp and os.path.exists(fp):
            if is_image:
                try:
                    image_obj = Image.open(fp).convert("RGB")
                    # Text-only fallback: use caption if available
                    cap = captions.get(r["task_id"], "")
                    file_text = f"[Image description: {cap}]" if cap else "[IMAGE FILE — no caption available]"
                except Exception as e:
                    file_text = f"[Image load failed: {e}]"
            elif ext in _EXTRACTORS and _EXTRACTORS[ext] is not None:
                extractor = _EXTRACTORS[ext]
                file_text = extractor(fp)
            else:
                # Unknown type — try reading as text
                file_text = _extract_txt(fp)
        elif fn and fp and not os.path.exists(fp):
            file_text = f"[FILE MISSING: {fn}]"
            skipped += 1

        samples.append({
            "task_id": r["task_id"],
            "level": int(r["level"]),
            "file_name": fn,
            "file_type": ext or "none",
            "question": r["Question"],
            "answer": r["Final answer"],
            "file_text": file_text,
            "has_image": is_image,
            "image": image_obj,
        })

    print(f"[load] {len(samples)} samples loaded, {skipped} files missing", flush=True)
    return samples


# ─── answer scoring ────────────────────────────────────────────────────────────

def normalize_gaia(s: str) -> str:
    if s is None:
        return ""
    s = str(s).strip().lower()
    s = re.sub(r"^(the |a |an )", "", s)
    s = s.replace(",", "")
    s = re.sub(r"\s+", " ", s)
    s = s.rstrip(" .,;:!?'\"")
    return s


def gaia_score(pred: str, gold: str) -> bool:
    np_ = normalize_gaia(pred)
    ng  = normalize_gaia(gold)
    if not ng:
        return False
    if ng == np_ or ng in np_:
        return True
    try:
        return abs(float(np_.replace(",", "")) - float(ng.replace(",", ""))) < 0.01
    except ValueError:
        return False


def extract_answer(text: str) -> str:
    if not text:
        return ""
    m = re.search(r"Answer\s*:?\s*\*{0,2}([^\n*]+?)\*{0,2}\s*$",
                  text, re.IGNORECASE | re.MULTILINE)
    if m:
        return m.group(1).strip().rstrip(".")
    m = re.search(r"final\s*answer\s*:?\s*\*{0,2}([^\n*]+?)\*{0,2}",
                  text, re.IGNORECASE)
    if m:
        return m.group(1).strip().rstrip(".")
    boxed = re.findall(r"\\boxed\{([^{}]+)\}", text)
    if boxed:
        return boxed[-1].strip()
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return lines[-1] if lines else text.strip()


# ─── prompt builders ────────────────────────────────────────────────────────────

_TOK = None
def _tok():
    global _TOK
    if _TOK is None:
        from transformers import AutoTokenizer
        _TOK = AutoTokenizer.from_pretrained(TEXT_PATH, trust_remote_code=True)
    return _TOK


def _to_ids(user_text: str) -> list[int]:
    msgs = [{"role": "user", "content": user_text}]
    p = _tok().apply_chat_template(
        msgs, tokenize=False, add_generation_prompt=True, enable_thinking=False)
    return _tok().encode(p)


def build_user_text(s: dict, trunc_tokens: int) -> str:
    """Verifier prompt: truncated file text + question."""
    if s["file_text"]:
        toks = _tok().encode(s["file_text"])[:trunc_tokens]
        ctx  = _tok().decode(toks, skip_special_tokens=True)
        return PROMPT_WITH_CTX.format(ctx=ctx, question=s["question"])
    return PROMPT_NO_CTX.format(question=s["question"])


def build_aug_user_text(s: dict) -> str:
    """Drafter prompt: full file text + question."""
    if s["file_text"]:
        return PROMPT_WITH_CTX.format(ctx=s["file_text"], question=s["question"])
    return PROMPT_NO_CTX.format(question=s["question"])


# ─── save ────────────────────────────────────────────────────────────────────

def save_results(cfg, samples, outs, out_path, spec_metrics=None, args=None):
    n_ok = sum(o["ok"] for o in outs)
    n   = len(outs)
    # Breakdown by level and file_type
    by_level: dict[int, dict] = {}
    by_type:  dict[str, dict] = {}
    for s, o in zip(samples[:n], outs):
        L = s["level"]; ft = s["file_type"] or "none"
        by_level.setdefault(L, {"n":0,"ok":0})
        by_type.setdefault(ft,  {"n":0,"ok":0})
        by_level[L]["n"] += 1; by_level[L]["ok"] += int(o["ok"])
        by_type[ft]["n"]  += 1; by_type[ft]["ok"]  += int(o["ok"])

    rec = {
        "cfg": cfg, "n_total": n, "n_correct": n_ok,
        "accuracy_pct": 100.0 * n_ok / n if n else 0.0,
        "by_level":    {k: {**v, "pct": 100*v["ok"]/v["n"]} for k,v in sorted(by_level.items())},
        "by_file_type":{k: {**v, "pct": 100*v["ok"]/v["n"]} for k,v in sorted(by_type.items())},
        "spec_metrics": spec_metrics,
        "args": vars(args) if args else {},
        "results": [
            {"task_id": s["task_id"], "level": s["level"],
             "file_type": s["file_type"], "question": s["question"],
             "gold": s["answer"], "pred": o["pred"], "ans": o.get("ans",""), "ok": o["ok"]}
            for s, o in zip(samples[:n], outs)
        ],
    }
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w") as f:
        json.dump(rec, f, indent=2, default=str)
    print(f"\n[saved] {cfg} → {out_path}  acc={rec['accuracy_pct']:.1f}% ({n_ok}/{n})", flush=True)
    for L, v in sorted(by_level.items()):
        print(f"  L{L}: {v['ok']}/{v['n']} = {100*v['ok']/v['n']:.1f}%", flush=True)


# ─── cfg: floor (Qwen3-32B alone) ─────────────────────────────────────────────

def cfg_floor(samples, out_path, args):
    print(f"\n=== floor: Qwen3-32B alone (trunc={args.trunc_tokens}) ===", flush=True)
    from vllm import LLM, SamplingParams
    llm = LLM(model=TEXT_PATH, dtype="bfloat16", trust_remote_code=True,
              max_model_len=args.max_model_len, gpu_memory_utilization=0.85,
              enforce_eager=True)
    sp = SamplingParams(temperature=0, max_tokens=args.max_new)
    outs, t0 = [], time.perf_counter()
    for i, s in enumerate(samples):
        try:
            prompt = _tok().apply_chat_template(
                [{"role": "user", "content": build_user_text(s, args.trunc_tokens)}],
                tokenize=False, add_generation_prompt=True, enable_thinking=False)
            r = llm.generate([{"prompt": prompt}], sp, use_tqdm=False)
            ans  = r[0].outputs[0].text.strip()
            pred = extract_answer(ans)
            ok   = gaia_score(pred, s["answer"])
            outs.append({"ans": ans, "pred": pred, "ok": ok})
            print(f"  [{i+1}/{len(samples)}] L{s['level']} {s['file_type']:<6} ok={ok}  "
                  f"pred={pred[:30]!r}  gold={s['answer']!r}", flush=True)
        except Exception as e:
            print(f"  [{i+1}] FAIL: {type(e).__name__}: {str(e)[:200]}", flush=True)
            outs.append({"ans": f"FAIL:{type(e).__name__}", "pred": "", "ok": False})
    print(f"[floor] elapsed={time.perf_counter()-t0:.1f}s", flush=True)
    save_results("floor", samples, outs, out_path, args=args)


# ─── cfg: aug (Qwen3-4B drafter alone) ────────────────────────────────────────

def cfg_aug(samples, out_path, args):
    print(f"\n=== aug: Qwen3-4B alone (full context) ===", flush=True)
    from vllm import LLM, SamplingParams
    from transformers import AutoTokenizer
    tok4 = AutoTokenizer.from_pretrained(DRAFTER_PATH, trust_remote_code=True)
    llm = LLM(model=DRAFTER_PATH, dtype="bfloat16", trust_remote_code=True,
              max_model_len=args.max_model_len, gpu_memory_utilization=0.85,
              enforce_eager=True)
    sp = SamplingParams(temperature=0, max_tokens=args.max_new)
    outs, t0 = [], time.perf_counter()
    for i, s in enumerate(samples):
        try:
            user_text = build_aug_user_text(s)
            prompt = tok4.apply_chat_template(
                [{"role": "user", "content": user_text}],
                tokenize=False, add_generation_prompt=True, enable_thinking=False)
            r = llm.generate([{"prompt": prompt}], sp, use_tqdm=False)
            ans  = r[0].outputs[0].text.strip()
            pred = extract_answer(ans)
            ok   = gaia_score(pred, s["answer"])
            outs.append({"ans": ans, "pred": pred, "ok": ok})
            print(f"  [{i+1}/{len(samples)}] L{s['level']} {s['file_type']:<6} ok={ok}  "
                  f"pred={pred[:30]!r}  gold={s['answer']!r}", flush=True)
        except Exception as e:
            print(f"  [{i+1}] FAIL: {type(e).__name__}: {str(e)[:200]}", flush=True)
            outs.append({"ans": f"FAIL:{type(e).__name__}", "pred": "", "ok": False})
    print(f"[aug] elapsed={time.perf_counter()-t0:.1f}s", flush=True)
    save_results("aug", samples, outs, out_path, args=args)


# ─── cfg: ss (SpecSteer text drafter) ─────────────────────────────────────────

def cfg_ss(samples, out_path, args):
    print(f"\n=== ss: SpecSteer K={args.K} (text drafter={DRAFTER_PATH}) ===", flush=True)
    from vllm import LLM, SamplingParams
    import vllm.v1.spec_decode.metrics as _sdm

    os.environ["ASYMSPEC_METHOD"] = args.asym_method
    os.environ["ASYMSPEC_DELTA_SRC"] = "ours"

    SPEC = {"drafts": 0, "draft_tokens": 0, "accepted_tokens": 0, "per_pos": None}
    _orig = _sdm.SpecDecodingStats.observe_draft
    def _capture(self, num_draft_tokens, num_accepted_tokens):
        SPEC["drafts"] += 1
        SPEC["draft_tokens"] += num_draft_tokens
        SPEC["accepted_tokens"] += num_accepted_tokens
        if SPEC["per_pos"] is None:
            SPEC["per_pos"] = [0] * self.num_spec_tokens
        for j in range(num_accepted_tokens):
            if j < len(SPEC["per_pos"]):
                SPEC["per_pos"][j] += 1
        return _orig(self, num_draft_tokens, num_accepted_tokens)
    _sdm.SpecDecodingStats.observe_draft = _capture

    llm = LLM(model=TEXT_PATH, dtype="bfloat16", trust_remote_code=True,
              max_model_len=args.max_model_len, gpu_memory_utilization=0.85,
              enforce_eager=True,
              disable_log_stats=False,
              speculative_config={
                  "method": "specsteer", "model": DRAFTER_PATH,
                  "num_speculative_tokens": args.K,
                  "specsteer_beta": args.beta, "specsteer_gamma": args.gamma,
              })
    print("[ss] LLM init OK", flush=True)
    sp_base = dict(temperature=0, max_tokens=args.max_new)
    outs, t0 = [], time.perf_counter()
    MAX_AUG = args.max_model_len - args.max_new - 32

    for i, s in enumerate(samples):
        try:
            main_text = build_user_text(s, args.trunc_tokens)
            aug_text  = build_aug_user_text(s)
            main_prompt = _tok().apply_chat_template(
                [{"role": "user", "content": main_text}],
                tokenize=False, add_generation_prompt=True, enable_thinking=False)
            aug_ids = _to_ids(aug_text)
            if len(aug_ids) > MAX_AUG:
                # truncate aug to fit; drafter still gets more than verifier
                aug_ids = aug_ids[:MAX_AUG]
            sp_kw = dict(sp_base, extra_args={"specsteer_aug_prompt_ids": aug_ids})
            r = llm.generate([{"prompt": main_prompt}],
                             SamplingParams(**sp_kw), use_tqdm=False)
            ans  = r[0].outputs[0].text.strip()
            pred = extract_answer(ans)
            ok   = gaia_score(pred, s["answer"])
            outs.append({"ans": ans, "pred": pred, "ok": ok})
            print(f"  [{i+1}/{len(samples)}] L{s['level']} {s['file_type']:<6} ok={ok}  "
                  f"pred={pred[:30]!r}  gold={s['answer']!r}", flush=True)
        except Exception as e:
            print(f"  [{i+1}] FAIL: {type(e).__name__}: {str(e)[:200]}", flush=True)
            outs.append({"ans": f"FAIL:{type(e).__name__}", "pred": "", "ok": False})

    el = time.perf_counter() - t0
    metrics = {
        "K": args.K, "beta": args.beta, "gamma": args.gamma,
        "asym_method": args.asym_method, "elapsed": el,
        "AR": SPEC["accepted_tokens"] / SPEC["draft_tokens"] if SPEC["draft_tokens"] else 0,
        "MAL": 1 + SPEC["accepted_tokens"] / SPEC["drafts"] if SPEC["drafts"] else 1,
        "per_position_acceptance_rate":
            [c / SPEC["drafts"] for c in (SPEC["per_pos"] or [])],
    }
    print(f"[SS metrics] AR={metrics['AR']:.3f} MAL={metrics['MAL']:.3f}", flush=True)
    save_results(f"ss_K{args.K}", samples, outs, out_path, spec_metrics=metrics, args=args)


# ─── cfg: ss_vl (VL drafter for image samples, text drafter for others) ───────

def cfg_ss_vl(samples, out_path, args):
    print(f"\n=== ss_vl: SpecSteer VL drafter K={args.K} ===", flush=True)
    from vllm import LLM, SamplingParams
    from transformers import AutoProcessor
    import vllm.v1.spec_decode.metrics as _sdm

    os.environ["ASYMSPEC_METHOD"] = args.asym_method
    os.environ["ASYMSPEC_DELTA_SRC"] = "ours"

    SPEC = {"drafts": 0, "draft_tokens": 0, "accepted_tokens": 0, "per_pos": None}
    _orig = _sdm.SpecDecodingStats.observe_draft
    def _capture(self, num_draft_tokens, num_accepted_tokens):
        SPEC["drafts"] += 1; SPEC["draft_tokens"] += num_draft_tokens
        SPEC["accepted_tokens"] += num_accepted_tokens
        if SPEC["per_pos"] is None:
            SPEC["per_pos"] = [0] * self.num_spec_tokens
        for j in range(num_accepted_tokens):
            if j < len(SPEC["per_pos"]): SPEC["per_pos"][j] += 1
        return _orig(self, num_draft_tokens, num_accepted_tokens)
    _sdm.SpecDecodingStats.observe_draft = _capture

    proc = AutoProcessor.from_pretrained(VL_PATH, trust_remote_code=True)
    llm = LLM(model=TEXT_PATH, dtype="bfloat16", trust_remote_code=True,
              max_model_len=args.max_model_len, gpu_memory_utilization=0.85,
              enforce_eager=True,
              disable_log_stats=False,
              speculative_config={
                  "method": "specsteer", "model": VL_PATH,
                  "num_speculative_tokens": args.K,
                  "specsteer_beta": args.beta, "specsteer_gamma": args.gamma,
              })
    print("[ss_vl] LLM init OK", flush=True)
    sp_base = dict(temperature=0, max_tokens=args.max_new)
    MAX_AUG = args.max_model_len - args.max_new - 32
    outs, t0 = [], time.perf_counter()

    for i, s in enumerate(samples):
        try:
            main_text = build_user_text(s, args.trunc_tokens)
            main_prompt = _tok().apply_chat_template(
                [{"role": "user", "content": main_text}],
                tokenize=False, add_generation_prompt=True, enable_thinking=False)

            if s["has_image"] and s["image"] is not None:
                # VL drafter: image + caption + question
                aug_user_text = build_aug_user_text(s)
                vl_msgs = [{"role": "user", "content": [
                    {"type": "image"}, {"type": "text", "text": aug_user_text}]}]
                aug_text = proc.apply_chat_template(
                    vl_msgs, tokenize=False, add_generation_prompt=True,
                    enable_thinking=False)
                po = proc(text=[aug_text], images=[s["image"]], return_tensors="pt",
                          padding=False)
                aug_ids = po["input_ids"][0].tolist()
                if len(aug_ids) > MAX_AUG:
                    aug_ids = aug_ids[:MAX_AUG]
                sp_kw = dict(sp_base, extra_args={
                    "specsteer_aug_prompt_ids": aug_ids,
                    "specsteer_aug_pixel_values": po["pixel_values"],
                    "specsteer_aug_image_grid_thw": po["image_grid_thw"],
                })
            else:
                # Text-only path: VL drafter reads text (no image)
                aug_text  = build_aug_user_text(s)
                vl_msgs = [{"role": "user", "content": aug_text}]
                aug_text_fmt = proc.apply_chat_template(
                    vl_msgs, tokenize=False, add_generation_prompt=True,
                    enable_thinking=False)
                aug_ids = proc(text=[aug_text_fmt], return_tensors="pt",
                               padding=False)["input_ids"][0].tolist()
                if len(aug_ids) > MAX_AUG:
                    aug_ids = aug_ids[:MAX_AUG]
                sp_kw = dict(sp_base, extra_args={"specsteer_aug_prompt_ids": aug_ids})

            r = llm.generate([{"prompt": main_prompt}],
                             SamplingParams(**sp_kw), use_tqdm=False)
            ans  = r[0].outputs[0].text.strip()
            pred = extract_answer(ans)
            ok   = gaia_score(pred, s["answer"])
            outs.append({"ans": ans, "pred": pred, "ok": ok})
            drafter_tag = "VL" if s["has_image"] else "txt"
            print(f"  [{i+1}/{len(samples)}] L{s['level']} {s['file_type']:<6} [{drafter_tag}] "
                  f"ok={ok}  pred={pred[:30]!r}  gold={s['answer']!r}", flush=True)
        except Exception as e:
            print(f"  [{i+1}] FAIL: {type(e).__name__}: {str(e)[:200]}", flush=True)
            outs.append({"ans": f"FAIL:{type(e).__name__}", "pred": "", "ok": False})

    el = time.perf_counter() - t0
    metrics = {
        "K": args.K, "beta": args.beta, "gamma": args.gamma,
        "asym_method": args.asym_method, "elapsed": el,
        "AR": SPEC["accepted_tokens"] / SPEC["draft_tokens"] if SPEC["draft_tokens"] else 0,
        "MAL": 1 + SPEC["accepted_tokens"] / SPEC["drafts"] if SPEC["drafts"] else 1,
        "per_position_acceptance_rate":
            [c / SPEC["drafts"] for c in (SPEC["per_pos"] or [])],
    }
    print(f"[SS_VL metrics] AR={metrics['AR']:.3f} MAL={metrics['MAL']:.3f}", flush=True)
    save_results(f"ss_vl_K{args.K}", samples, outs, out_path, spec_metrics=metrics, args=args)


# ─── main ──────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--cfg", required=True, choices=["floor", "aug", "ss", "ss_vl"])
    ap.add_argument("--K",    type=int,   default=2)
    ap.add_argument("--beta", type=float, default=1.0)
    ap.add_argument("--gamma",type=float, default=0.5)
    ap.add_argument("--asym_method", default="jsd")
    ap.add_argument("--trunc_tokens", type=int, default=2000,
                    help="Max tokens for verifier context (file content truncation)")
    ap.add_argument("--max_model_len", type=int, default=32768)
    ap.add_argument("--max_new",  type=int, default=1024)
    ap.add_argument("--out_dir",  default=OUT_DIR)
    ap.add_argument("--tag",      default="")
    ap.add_argument("--skip_image", action="store_true",
                    help="Exclude the 10 image samples (for text-only analysis)")
    ap.add_argument("--only_image", action="store_true",
                    help="Run only the 10 image samples")
    ap.add_argument("--only_with_file", action="store_true",
                    help="Run only the 38 file-bearing samples (skip 127 web-only)")
    args = ap.parse_args()

    samples = load_all_samples()
    if args.only_with_file:
        samples = [s for s in samples if s["file_name"]]
        print(f"[filter] file-bearing only → n={len(samples)}", flush=True)
    if args.skip_image:
        samples = [s for s in samples if not s["has_image"]]
        print(f"[filter] image samples excluded → n={len(samples)}", flush=True)
    elif args.only_image:
        samples = [s for s in samples if s["has_image"]]
        print(f"[filter] image-only → n={len(samples)}", flush=True)

    tag = args.tag or args.cfg
    suffix = f"_K{args.K}" if args.cfg in ("ss","ss_vl") else ""
    out_path = os.path.join(args.out_dir, f"{tag}{suffix}_n{len(samples)}.json")

    if args.cfg == "floor":
        cfg_floor(samples, out_path, args)
    elif args.cfg == "aug":
        cfg_aug(samples, out_path, args)
    elif args.cfg == "ss":
        cfg_ss(samples, out_path, args)
    else:
        cfg_ss_vl(samples, out_path, args)


if __name__ == "__main__":
    main()
